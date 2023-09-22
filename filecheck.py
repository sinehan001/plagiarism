import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation

def extract_data(file_path):
    file_extension = file_path.split('.')[-1].lower()

    if file_extension in ('doc', 'docx'):
        doc = Document(file_path)
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        return text

    elif file_extension == 'pdf':
        pdf_file = open(file_path, 'rb')
        pdf_reader = PyPDF2.PdfFileReader(pdf_file)
        text = ''
        for page_num in range(pdf_reader.numPages):
            page = pdf_reader.getPage(page_num)
            text += page.extractText()
        pdf_file.close()
        return text

    elif file_extension in ('xlsx', 'xls'):
        workbook = openpyxl.load_workbook(file_path)
        sheet = workbook.active
        data = []
        for row in sheet.iter_rows(values_only=True):
            data.append(row)
        return data

    elif file_extension == 'pptx':
        presentation = Presentation(file_path)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        return text

    else:
        return "Unsupported file format"

# Example usage
file_path = 'example.docx'  # Replace with the path to your file
extracted_data = extract_data(file_path)
print(extracted_data)
